"""
Script untuk membuat Buku Laporan Kegiatan Reses Anggota Fraksi (PPTX)
Masa Sidang V Tahun Sidang 2025-2026 - Ir. Andreas Eddy Susetyo, MM (A-220)
Periode Reses: 22 Juli - 13 Agustus 2026
Daerah Pemilihan Jawa Timur V (Kota Malang, Kabupaten Malang, Kota Batu)

PENTING: Script ini TIDAK membuat layout sendiri. Seluruh slide diedit langsung di atas
salinan template baku FINAL_BUKU_LAPORAN_RESES_Andreas_Eddy_Susetyo_April_Mei_2026.pptx.
Hanya teks, angka, dan foto yang diganti; posisi, ukuran, warna, dan format asli template
tidak diubah sama sekali.
"""

import io
import os
import shutil
import copy

from pptx import Presentation
from pptx.util import Emu
from PIL import Image

from create_reses_docx import KEGIATAN, DAPIL

# ============================================================
# CONFIG
# ============================================================
TEMPLATE_FILE = r"D:\DATA PAPA\LAPORAN\RESES\APRIL_MEI 2026\FINAL_BUKU_LAPORAN_RESES_Andreas_Eddy_Susetyo_April_Mei_2026.pptx"
OUTPUT_DIR = r"D:\DATA PAPA\LAPORAN\RESES\JULI_AGUSTUS 2026"
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "FINAL_BUKU_LAPORAN_RESES_Andreas_Eddy_Susetyo_Juli_Agustus_2026.pptx")
FOTO_DIR = r"D:\DATA PAPA\LAPORAN\RESES\JULI_AGUSTUS 2026"

TANGGAL_LAPORAN = "15 Agustus 2026"

# ============================================================
# DATA TAMBAHAN PER KEGIATAN (untuk Lembar Laporan Pelaksanaan)
# ============================================================
EXTRA_DATA = {
    1: {"kabkota": "Kabupaten Malang", "jumlah": 60,
        "dihadiri": "Dihadiri pengurus, anggota koperasi, dan simpatisan partai",
        "tanggapan": [
            "Akan diperjuangkan skema pembiayaan bunga rendah bagi koperasi melalui koordinasi dengan Kementerian Koperasi dan UKM serta lembaga perbankan mitra pemerintah.",
            "Akan diusulkan program pelatihan manajemen keuangan dan digitalisasi koperasi melalui kerja sama dengan Dinas Koperasi Kabupaten Malang dan Kementerian Koperasi dan UKM.",
            "Akan diadvokasikan insentif fiskal bagi koperasi aktif kepada Kementerian Keuangan dan Kementerian Koperasi dan UKM sebagai bagian pembahasan kebijakan anggaran.",
        ]},
    2: {"kabkota": "Kota Malang", "jumlah": 45,
        "dihadiri": "Dihadiri pelaku industri jasa keuangan dan simpatisan partai",
        "tanggapan": [
            "Akan dikoordinasikan program literasi keuangan masyarakat bersama OJK dan industri jasa keuangan guna menekan maraknya pinjaman daring ilegal.",
            "Akan diperjuangkan insentif kebijakan bagi lembaga keuangan formal untuk memperluas kredit UMKM melalui pembahasan bersama OJK dan Kementerian Keuangan.",
            "Akan didorong penguatan pengawasan OJK terhadap praktik fintech lending guna melindungi konsumen dari suku bunga tidak wajar.",
        ]},
    3: {"kabkota": "Kota Malang", "jumlah": 150,
        "dihadiri": "Dihadiri orang tua, siswa, dan simpatisan partai",
        "tanggapan": [
            "Akan diusulkan peningkatan anggaran sarana dan prasarana praktik SMK melalui koordinasi dengan Kementerian Pendidikan dan Pemerintah Kota Malang.",
            "Akan diperjuangkan perluasan program beasiswa dan bantuan biaya pendidikan bagi siswa kurang mampu melalui Komisi X dan Kementerian Pendidikan.",
            "Akan difasilitasi penguatan kerja sama sekolah dengan dunia usaha dan industri untuk memperluas peluang magang dan penyerapan lulusan SMK.",
        ]},
    4: {"kabkota": "Kota Batu", "jumlah": 250,
        "dihadiri": "Dihadiri tokoh agama dan umat Katolik setempat",
        "tanggapan": [
            "Akan dikoordinasikan dukungan pemerintah daerah untuk pemeliharaan dan pengembangan fasilitas sarana ibadah serta kegiatan sosial keagamaan.",
            "Akan didorong penguatan program kerukunan lintas iman yang melibatkan generasi muda melalui koordinasi dengan Kementerian Agama dan Pemerintah Kota Batu.",
            "Akan diadvokasikan dukungan terhadap program sosial keagamaan seperti layanan kesehatan dan pendidikan bagi jemaat kurang mampu.",
        ]},
    5: {"kabkota": "Kota Malang", "jumlah": 30,
        "dihadiri": "Dihadiri jajaran Perum Perhutani KPH Malang",
        "tanggapan": [
            "Akan diperjuangkan perluasan program perhutanan sosial agar lebih banyak masyarakat memperoleh akses legal pemanfaatan lahan hutan.",
            "Akan didorong peningkatan koordinasi antara Perhutani, pemerintah daerah, dan masyarakat dalam pencegahan kebakaran hutan dan lahan.",
            "Akan diusulkan program kemitraan ekonomi kehutanan bagi kelompok tani hutan melalui koordinasi dengan Kementerian Lingkungan Hidup dan Kehutanan.",
        ]},
    6: {"kabkota": "Kota Malang", "jumlah": 120,
        "dihadiri": "Dihadiri Ketua dan jajaran pengurus struktural DPC Kota Malang",
        "tanggapan": [
            "Akan ditindaklanjuti penguatan konsolidasi struktur partai hingga tingkat ranting melalui rapat koordinasi internal DPC Kota Malang.",
            "Akan didorong sinergi program kerja partai dengan program pemerintah Kota Malang agar aspirasi konstituen dapat ditindaklanjuti secara konkret.",
            "Akan diperkuat program kaderisasi dan pelatihan politik bagi kader muda partai di Kota Malang.",
        ]},
    7: {"kabkota": "Kabupaten Malang", "jumlah": 40,
        "dihadiri": "Dihadiri pelaku BPR dan simpatisan partai",
        "tanggapan": [
            "Akan diadvokasikan relaksasi kebijakan permodalan BPR kepada OJK agar mampu bersaing dengan lembaga keuangan digital dan perbankan nasional.",
            "Akan diusulkan penyederhanaan regulasi kepatuhan BPR kepada OJK agar tidak membebani operasional lembaga keuangan skala kecil.",
            "Akan difasilitasi perluasan kerja sama BPR dengan program pembiayaan UMKM dan sektor pertanian bersama Pemerintah Kabupaten Malang.",
        ]},
    8: {"kabkota": "Kabupaten Malang", "jumlah": 55,
        "dihadiri": "Dihadiri pengurus koperasi simpan pinjam dan simpatisan partai",
        "tanggapan": [
            "Akan diusulkan program pendampingan teknis manajemen risiko kredit bagi pengurus koperasi simpan pinjam melalui Dinas Koperasi.",
            "Akan didorong penegakan hukum dan pengawasan terhadap praktik simpan pinjam ilegal berkedok koperasi.",
            "Akan diperjuangkan akses permodalan tambahan dari lembaga keuangan pemerintah untuk memperkuat kapasitas koperasi.",
        ]},
    9: {"kabkota": "Kota Malang", "jumlah": 300,
        "dihadiri": "Dihadiri tokoh agama dan jemaat gereja",
        "tanggapan": [
            "Akan diusulkan dukungan pemerintah untuk pemeliharaan bangunan gereja bersejarah sebagai bagian dari cagar budaya kota.",
            "Akan terus didorong penguatan program toleransi dan kerukunan umat beragama di lingkungan perkotaan.",
            "Akan dikoordinasikan dukungan keamanan dan ketertiban bagi penyelenggaraan kegiatan keagamaan berskala besar bersama aparat terkait.",
        ]},
    10: {"kabkota": "Kabupaten Malang", "jumlah": 90,
        "dihadiri": "Dihadiri tokoh masyarakat dan warga setempat",
        "tanggapan": [
            "Akan diusulkan perbaikan infrastruktur jalan dan drainase lingkungan kepada Dinas Pekerjaan Umum Kabupaten Malang.",
            "Akan didorong optimalisasi layanan kesehatan dasar di Puskesmas setempat melalui koordinasi dengan Dinas Kesehatan.",
            "Akan dikawal penyaluran bantuan sosial yang lebih tepat sasaran bagi warga terdampak pelemahan daya beli.",
        ]},
    11: {"kabkota": "Kabupaten Malang", "jumlah": 70,
        "dihadiri": "Dihadiri perempuan pelaku UMKM dan simpatisan partai",
        "tanggapan": [
            "Akan difasilitasi pelatihan digital marketing dan akses pembiayaan mikro khusus bagi perempuan pelaku usaha.",
            "Akan diusulkan program pendampingan perizinan usaha seperti NIB dan sertifikasi halal bagi UMKM perempuan.",
            "Akan diperjuangkan fasilitasi pemasaran produk UMKM perempuan melalui pameran dan kemitraan dengan ritel modern.",
        ]},
    12: {"kabkota": "Kabupaten Malang", "jumlah": 65,
        "dihadiri": "Dihadiri tokoh masyarakat dan warga Desa Druju",
        "tanggapan": [
            "Akan diusulkan perbaikan akses jalan desa kepada Pemerintah Kabupaten Malang guna melancarkan distribusi hasil pertanian dan perkebunan.",
            "Akan dikoordinasikan perluasan jaringan listrik dan telekomunikasi di wilayah pedesaan bersama PLN dan operator telekomunikasi.",
            "Akan diperjuangkan bantuan bibit unggul dan pupuk bersubsidi melalui Dinas Pertanian Kabupaten Malang.",
        ]},
    13: {"kabkota": "Kabupaten Malang", "jumlah": 100,
        "dihadiri": "Dihadiri Ketua dan jajaran pengurus struktural PAC Kabupaten Malang",
        "tanggapan": [
            "Akan ditindaklanjuti penguatan koordinasi antara PAC dan Ranting dalam menjaring aspirasi masyarakat tingkat kecamatan.",
            "Akan didorong dukungan program pemberdayaan ekonomi berbasis potensi wisata lokal Kecamatan Tumpang bersama Pemerintah Kabupaten Malang.",
            "Akan diperkuat program kaderisasi politik generasi muda sebagai bagian regenerasi kepemimpinan partai di daerah.",
        ]},
    14: {"kabkota": "Kabupaten Malang", "jumlah": 85,
        "dihadiri": "Dihadiri petani, nelayan, dan simpatisan partai",
        "tanggapan": [
            "Akan diadvokasikan intervensi pasar oleh pemerintah guna menjaga stabilitas harga hasil pertanian dan tangkapan nelayan.",
            "Akan diusulkan bantuan sarana produksi seperti alat tangkap ikan dan pupuk bersubsidi bagi petani dan nelayan pesisir selatan.",
            "Akan diperjuangkan pembangunan tempat pelelangan ikan dan akses jalan menuju kawasan pesisir melalui Kementerian Kelautan dan Perikanan.",
        ]},
    15: {"kabkota": "Kabupaten Malang", "jumlah": 75,
        "dihadiri": "Dihadiri tokoh masyarakat dan warga Desa Sambigede",
        "tanggapan": [
            "Akan diusulkan perbaikan sistem irigasi pertanian kepada Dinas Pengairan dan Pertanian Kabupaten Malang.",
            "Akan dikoordinasikan peningkatan akses air bersih bagi warga bersama PDAM dan Pemerintah Kabupaten Malang.",
            "Akan diperjuangkan program padat karya untuk menyerap tenaga kerja lokal di sektor infrastruktur desa.",
        ]},
    16: {"kabkota": "Kabupaten Malang", "jumlah": 60,
        "dihadiri": "Dihadiri pelaku UMKM dan simpatisan partai",
        "tanggapan": [
            "Akan difasilitasi kemudahan akses Kredit Usaha Rakyat bagi pelaku UMKM di wilayah pedesaan bersama lembaga perbankan penyalur KUR.",
            "Akan diusulkan pelatihan pengemasan dan standardisasi produk UMKM melalui Dinas Koperasi dan UKM Kabupaten Malang.",
            "Akan diperjuangkan perbaikan akses jalan produksi guna memperlancar distribusi hasil UMKM ke luar wilayah kecamatan.",
        ]},
    17: {"kabkota": "Kota Batu", "jumlah": 110,
        "dihadiri": "Dihadiri Ketua dan jajaran pengurus struktural DPC Kota Batu",
        "tanggapan": [
            "Akan ditindaklanjuti penguatan konsolidasi kader di tingkat kelurahan guna memperkuat basis dukungan partai di Kota Batu.",
            "Akan didorong sinkronisasi program partai dengan agenda pembangunan pariwisata Kota Batu bersama Pemerintah Kota Batu.",
            "Akan diperkuat peningkatan kapasitas kader dalam mengawasi pelaksanaan program pemerintah daerah.",
        ]},
    18: {"kabkota": "Kabupaten Malang", "jumlah": 50,
        "dihadiri": "Dihadiri warga dan simpatisan partai",
        "tanggapan": [
            "Akan diusulkan dukungan anggaran perbaikan fasilitas umum desa seperti saluran air dan jalan lingkungan kepada Pemerintah Kabupaten Malang.",
            "Akan terus didorong penguatan program gotong royong sebagai bagian pemberdayaan masyarakat berbasis kearifan lokal.",
            "Akan dikoordinasikan program pengelolaan sampah dan kebersihan lingkungan bersama Dinas Lingkungan Hidup Kabupaten Malang.",
        ]},
    19: {"kabkota": "Kabupaten Malang", "jumlah": 55,
        "dihadiri": "Dihadiri tokoh masyarakat dan warga Desa Ngajum",
        "tanggapan": [
            "Akan diusulkan perbaikan infrastruktur jalan desa yang menjadi akses utama distribusi hasil pertanian.",
            "Akan diperjuangkan bantuan alat mesin pertanian dan bibit unggul melalui Dinas Pertanian Kabupaten Malang.",
            "Akan dikoordinasikan penyediaan akses air bersih bagi warga Desa Ngajum pada musim kemarau bersama PDAM.",
        ]},
    20: {"kabkota": "Kabupaten Malang", "jumlah": 400,
        "dihadiri": "Dihadiri pegiat seni budaya dan masyarakat umum",
        "tanggapan": [
            "Akan diusulkan dukungan pemerintah daerah terhadap pelestarian seni budaya jaranan sebagai warisan budaya lokal.",
            "Akan diperjuangkan fasilitasi panggung dan anggaran penyelenggaraan festival budaya secara berkala bersama Dinas Kebudayaan dan Pariwisata.",
            "Akan didorong program regenerasi pelaku seni jaranan melalui pelatihan bagi generasi muda agar tradisi ini tetap lestari.",
        ]},
}

# ============================================================
# TEKS PENGANTAR RESES - 7 paragraf isi (4 di slide Pengantar #1, 3 di slide Pengantar #2)
# menggantikan paragraf isi pada template, judul & paragraf kosong pengisi tidak disentuh.
# ============================================================
PENGANTAR_PARAS = [
    # -- slide Pengantar Reses ke-1 (4 paragraf) --
    "Sesuai dengan amanat konstitusi dan peraturan perundang-undangan, Anggota Dewan "
    "Perwakilan Rakyat Republik Indonesia telah melaksanakan kegiatan reses pada periode "
    "22 Juli hingga 13 Agustus 2026 di Daerah Pemilihan Jawa Timur V yang meliputi Kota "
    "Malang, Kabupaten Malang, dan Kota Batu, pada Masa Sidang V Tahun Sidang 2025-2026. "
    "Kegiatan reses merupakan wujud nyata pelaksanaan fungsi representasi Anggota DPR RI "
    "dalam menjaga komunikasi dan kedekatan langsung dengan konstituen di daerah "
    "pemilihan, sekaligus memenuhi kewajiban konstitusional sebagai wakil rakyat yang "
    "terpilih melalui pemilihan umum.",

    "Salah satu fokus utama pelaksanaan reses kali ini adalah memantau dan mencari "
    "masukan langsung terkait kondisi perekonomian masyarakat di Malang Raya. Sebagai "
    "Anggota Komisi XI DPR RI yang membidangi keuangan dan perencanaan pembangunan, "
    "penting untuk melihat secara langsung kondisi riil di lapangan, mulai dari daya beli "
    "masyarakat, keberlangsungan usaha koperasi dan UMKM, kinerja lembaga keuangan mikro "
    "seperti BPR dan koperasi simpan pinjam, hingga tantangan yang dihadapi sektor "
    "pertanian, perikanan, dan kehutanan yang menjadi tulang punggung perekonomian di "
    "sebagian besar wilayah Kabupaten Malang.",

    "Selain aspek ekonomi, reses ini juga diarahkan untuk memantau dan menggali kondisi "
    "sosial masyarakat secara lebih menyeluruh. Dinamika sosial yang berkembang di tengah "
    "masyarakat, baik pada aspek pendidikan, kehidupan keagamaan, kegotongroyongan, "
    "maupun pelestarian seni dan budaya lokal, menjadi bagian penting yang tidak dapat "
    "dipisahkan dari upaya membangun kesejahteraan masyarakat secara utuh.",

    "Melalui rangkaian dialog dan kunjungan lapangan, reses ini juga bertujuan untuk "
    "menggali secara langsung berbagai problematika yang dihadapi masyarakat sehari-hari, "
    "seperti keterbatasan infrastruktur jalan dan irigasi, minimnya akses air bersih, "
    "keterbatasan sarana pendidikan dan kesehatan, hingga sulitnya akses permodalan bagi "
    "pelaku usaha kecil. Problematika tersebut diidentifikasi secara mendalam agar dapat "
    "dirumuskan solusi kebijakan yang tepat sasaran.",

    # -- slide Pengantar Reses ke-2 (3 paragraf) --
    "Momentum reses kali ini menjadi semakin strategis karena bersamaan dengan "
    "berlangsungnya pembahasan Rancangan Anggaran Pendapatan dan Belanja Negara (RAPBN) "
    "Tahun Anggaran 2027 di tingkat pusat. Sebagai Anggota Komisi XI yang turut membahas "
    "kebijakan fiskal dan anggaran negara, masukan yang diserap langsung dari masyarakat "
    "di daerah pemilihan menjadi bahan pengayaan yang sangat berharga untuk memastikan "
    "alokasi anggaran negara benar-benar responsif terhadap kebutuhan riil masyarakat di "
    "daerah.",

    "Laporan ini secara khusus memuat rangkaian kegiatan reses yang dilaksanakan pada "
    "periode 27 Juli hingga 4 Agustus 2026, terdiri atas 20 kegiatan yang tersebar di "
    "wilayah Kota Malang, Kabupaten Malang, dan Kota Batu. Rangkaian kegiatan tersebut "
    "mencakup berbagai elemen masyarakat, mulai dari pelaku ekonomi seperti koperasi, "
    "UMKM, BPR, dan industri jasa keuangan; kalangan pendidikan; tokoh dan umat lintas "
    "agama; jajaran struktur partai; kelompok tani dan nelayan; hingga pegiat seni dan "
    "budaya lokal.",

    "Pelaksanaan reses dilakukan melalui dialog interaktif, kunjungan lapangan, serta "
    "pertemuan formal dengan mitra kerja dan elemen masyarakat di daerah pemilihan, "
    "dengan pendekatan yang partisipatif dan terbuka agar aspirasi yang terserap "
    "benar-benar mencerminkan kebutuhan riil masyarakat. Laporan ini disusun sebagai "
    "bentuk pertanggungjawaban pelaksanaan tugas konstitusional sekaligus dokumentasi "
    "menyeluruh atas seluruh rangkaian kegiatan Kunjungan Kerja Reses Perseorangan, agar "
    "seluruh aspirasi yang telah diserap dapat ditindaklanjuti melalui forum rapat kerja "
    "Komisi XI, pembahasan RAPBN 2027, maupun koordinasi dengan instansi pemerintah "
    "terkait.",
]

# ============================================================
# TEKS PENUTUP - 3 paragraf isi (judul & paragraf kosong pengisi tidak disentuh)
# ============================================================
PENUTUP_PARAS = [
    "Kegiatan reses periode 22 Juli hingga 13 Agustus 2026 ini telah berhasil "
    "mengumpulkan berbagai aspirasi dan masukan dari masyarakat Daerah Pemilihan Jawa "
    "Timur V yang akan menjadi bahan kerja sebagai wakil rakyat. Rangkaian 20 kegiatan "
    "yang dilaksanakan pada 27 Juli hingga 4 Agustus 2026 berhasil menjangkau pelaku "
    "ekonomi, kalangan pendidikan, tokoh lintas agama, jajaran struktur partai, petani "
    "dan nelayan, hingga pegiat seni budaya di Kota Malang, Kabupaten Malang, dan Kota "
    "Batu. Seluruh aspirasi yang terkumpul akan didokumentasikan serta menjadi bahan "
    "dalam setiap perumusan kebijakan, khususnya dalam pembahasan RAPBN Tahun Anggaran "
    "2027.",

    "Temuan lapangan menunjukkan bahwa akses permodalan dengan syarat mudah dan bunga "
    "terjangkau masih menjadi kebutuhan sentral bagi pelaku usaha, sementara pemerataan "
    "infrastruktur dasar seperti jalan, irigasi, air bersih, dan sarana produksi "
    "pertanian maupun perikanan masih menjadi kebutuhan mendesak di berbagai desa. Di "
    "sisi lain, kohesi sosial dan pelestarian nilai budaya lokal tetap menjadi kekuatan "
    "utama masyarakat Malang Raya yang harus terus dijaga dan diperkuat.",

    "Sebagai tindak lanjut, Anggota DPR RI berkomitmen untuk memperjuangkan aspirasi "
    "masyarakat melalui berbagai mekanisme parlemen, seperti rapat dengar pendapat, "
    "kunjungan kerja, dan pembahasan kebijakan anggaran di Komisi XI. Komunikasi "
    "berkelanjutan dengan masyarakat akan terus dijaga untuk memastikan bahwa suara "
    "rakyat benar-benar terdengar dan diperjuangkan di tingkat nasional, sehingga reses "
    "menjadi bagian integral dari upaya membangun demokrasi yang lebih partisipatif, "
    "responsif, dan berpihak pada rakyat.",
]


# ============================================================
# HELPER FUNCTIONS - edit teks tanpa mengubah format asli template
# ============================================================

def set_paragraph_text(paragraph, text):
    """Ganti teks sebuah paragraf, pertahankan format run pertama, hapus run sisanya."""
    runs = paragraph.runs
    if not runs:
        run = paragraph.add_run()
    else:
        run = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    run.text = text


def set_body_paragraphs(text_frame, start_idx, texts):
    """Ganti teks paragraf mulai indeks start_idx sebanyak len(texts), sisanya (judul,
    paragraf kosong pengisi) tidak disentuh sama sekali."""
    paragraphs = text_frame.paragraphs
    for offset, text in enumerate(texts):
        set_paragraph_text(paragraphs[start_idx + offset], text)


def set_cell_first_paragraph(cell, text):
    """Set paragraf pertama, lalu KOSONGKAN sisa paragraf agar tidak ada teks lama
    dari template yang tertinggal (beberapa sel template ternyata punya baris kedua
    berisi data asli, bukan baris kosong)."""
    paragraphs = cell.text_frame.paragraphs
    set_paragraph_text(paragraphs[0], text)
    for p in paragraphs[1:]:
        set_paragraph_text(p, "")


def set_cell_all_paragraphs(cell, texts):
    paragraphs = cell.text_frame.paragraphs
    for offset, text in enumerate(texts):
        set_paragraph_text(paragraphs[offset], text)
    for p in paragraphs[len(texts):]:
        set_paragraph_text(p, "")


def replace_picture(pic_shape, new_image_path):
    """Ganti isi gambar sebuah Picture shape TANPA mengubah posisi/ukuran/transform.
    Foto sumber di-crop (bukan di-stretch) supaya mengisi bingkai asli tanpa distorsi."""
    if not os.path.exists(new_image_path):
        print(f"  [SKIP] Foto tidak ditemukan: {new_image_path}")
        return
    rId = pic_shape._element.blip_rId
    image_part = pic_shape.part.related_part(rId)
    ext = image_part.ext.lower()

    target_ratio = pic_shape.width / pic_shape.height

    with Image.open(new_image_path) as im:
        im = im.convert("RGB")
        iw, ih = im.size
        src_ratio = iw / ih
        if src_ratio > target_ratio:
            new_w = int(ih * target_ratio)
            offset = max(0, (iw - new_w) // 2)
            im = im.crop((offset, 0, offset + new_w, ih))
        else:
            new_h = int(iw / target_ratio)
            offset = max(0, (ih - new_h) // 2)
            im = im.crop((0, offset, iw, offset + new_h))
        buf = io.BytesIO()
        if ext == "png":
            im.save(buf, format="PNG")
        else:
            im.save(buf, format="JPEG", quality=90)
        buf.seek(0)

    image_part._blob = buf.read()
    pic_shape.crop_left = 0
    pic_shape.crop_top = 0
    pic_shape.crop_right = 0
    pic_shape.crop_bottom = 0


def find_pictures(shapes, acc):
    for shape in shapes:
        if shape.shape_type == 13:  # PICTURE
            acc.append(shape)
        elif shape.shape_type == 6:  # GROUP
            find_pictures(shape.shapes, acc)


def find_content_pictures(slide):
    """Ambil foto konten kegiatan pada slide Lampiran Foto, tidak termasuk logo header.

    Logo header (DPR/Fraksi) berada dalam grup terpisah yang menempel di bagian atas
    slide (top-level top kecil, ~0.3 inci). Grup foto konten selalu terletak jauh di
    bawahnya (top-level top > 1.5 inci). Posisi top di dalam grup bersifat relatif
    terhadap grup masing-masing sehingga tidak bisa dibandingkan lintas grup - filter
    ini hanya membandingkan top-level top antar-grup, yang berada dalam satu ruang
    koordinat slide yang sama, lalu baru mengurutkan foto konten berdasarkan top lokal
    di dalam grup yang sama.
    """
    acc = []
    for shape in slide.shapes:
        if shape.shape_type == 13:
            acc.append(shape)
        elif shape.shape_type == 6:
            if shape.top is not None and shape.top < Emu(1371600):  # ~1.5 inch: grup header
                continue
            find_pictures(shape.shapes, acc)
    acc.sort(key=lambda p: p.top)
    return acc


def find_tables(shapes):
    return [sh.table for sh in shapes if sh.has_table]


# ============================================================
# MAIN
# ============================================================

def create_presentation():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    shutil.copyfile(TEMPLATE_FILE, OUTPUT_FILE)
    prs = Presentation(OUTPUT_FILE)
    slides = prs.slides

    # ------------------------------------------------------------
    # SLIDE 1-2: SAMPUL - ganti nomor Masa Sidang & tanggal reses saja
    # ------------------------------------------------------------
    for idx in (0, 1):
        s = slides[idx]
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "PERSIDANGAN IV" in run.text:
                        run.text = run.text.replace("PERSIDANGAN IV", "PERSIDANGAN V")
                    if "Masa Reses:" in run.text:
                        run.text = "Masa Reses: 22 Juli 2026 - 13 Agustus 2026"

    # ------------------------------------------------------------
    # SLIDE 3: DAFTAR ISI - generik, tidak diubah
    # ------------------------------------------------------------

    # ------------------------------------------------------------
    # SLIDE 4-5: PENGANTAR RESES - ganti isi paragraf, judul & padding dibiarkan
    # ------------------------------------------------------------
    s4 = slides[3]
    s5 = slides[4]
    ph4 = next(sh for sh in s4.shapes if sh.has_text_frame and sh.text_frame.text.strip())
    ph5 = next(sh for sh in s5.shapes if sh.has_text_frame and sh.text_frame.text.strip())
    set_body_paragraphs(ph4.text_frame, 1, PENGANTAR_PARAS[0:4])
    set_body_paragraphs(ph5.text_frame, 0, PENGANTAR_PARAS[4:7])

    # ------------------------------------------------------------
    # SLIDE 6-7: LEMBAR LAPORAN PELAKSANAAN - TABEL RINGKASAN
    # ------------------------------------------------------------
    keg_cursor = 0
    for slide_idx in (5, 6):
        s = slides[slide_idx]
        tables = find_tables(s.shapes)
        for tbl in tables:
            has_header = tbl.cell(0, 0).text.strip().lower() == "no"
            row_range = range(1, len(tbl.rows)) if has_header else range(0, len(tbl.rows))
            for r in row_range:
                keg = KEGIATAN[keg_cursor]
                extra = EXTRA_DATA[keg["no"]]
                set_cell_first_paragraph(tbl.cell(r, 0), str(keg["no"]))
                set_cell_first_paragraph(tbl.cell(r, 1), keg["judul"])
                set_cell_all_paragraphs(tbl.cell(r, 2), [
                    keg["hari_tanggal"],
                    f'Pukul: {keg["waktu"]}',
                    f'Tempat: {keg["tempat"]}',
                ])
                set_cell_first_paragraph(tbl.cell(r, 3), f'{extra["jumlah"]} orang')
                keg_cursor += 1

    # ------------------------------------------------------------
    # SLIDE 8-27: LEMBAR LAPORAN PELAKSANAAN - DETAIL PER KEGIATAN
    # ------------------------------------------------------------
    for i, keg in enumerate(KEGIATAN):
        s = slides[7 + i]
        extra = EXTRA_DATA[keg["no"]]
        tables = find_tables(s.shapes)
        id_table = next(t for t in tables if len(t.rows) == 6)
        isu_table, tanggapan_table = None, None
        for t in tables:
            if len(t.rows) == 1 and len(t.columns) == 1:
                label = t.cell(0, 0).text_frame.paragraphs[0].text.strip().lower()
                if label.startswith("isu"):
                    isu_table = t
                elif label.startswith("tanggapan"):
                    tanggapan_table = t

        set_cell_first_paragraph(id_table.cell(0, 1), ": Ir. Andreas Eddy Susetyo, MM")
        set_cell_first_paragraph(id_table.cell(1, 1), f': {extra["kabkota"]}')
        set_cell_first_paragraph(id_table.cell(2, 1), f': {keg["hari_tanggal"]}, {keg["waktu"]}')
        set_cell_first_paragraph(id_table.cell(3, 1), f': {keg["judul"]}')
        set_cell_first_paragraph(id_table.cell(4, 1), f': {extra["jumlah"]} orang')
        set_cell_first_paragraph(id_table.cell(5, 1), f': {extra["dihadiri"]}')

        isu_paras = isu_table.cell(0, 0).text_frame.paragraphs
        for offset, u in enumerate(keg["usulan"]):
            set_paragraph_text(isu_paras[1 + offset], f"{offset + 1})  {u}")
        for p in isu_paras[1 + len(keg["usulan"]):]:
            set_paragraph_text(p, "")

        tanggapan_paras = tanggapan_table.cell(0, 0).text_frame.paragraphs
        for offset, t in enumerate(extra["tanggapan"]):
            set_paragraph_text(tanggapan_paras[1 + offset], f"{offset + 1})  {t}")
        for p in tanggapan_paras[1 + len(extra["tanggapan"]):]:
            set_paragraph_text(p, "")

    # ------------------------------------------------------------
    # SLIDE 28: PENUTUP + tanda tangan
    # ------------------------------------------------------------
    s28 = slides[27]
    penutup_ph = next(sh for sh in s28.shapes
                      if sh.has_text_frame and sh.text_frame.text.strip().startswith("PENUTUP"))
    set_body_paragraphs(penutup_ph.text_frame, 1, PENUTUP_PARAS)

    for tbl in find_tables(s28.shapes):
        first_cell_text = tbl.cell(0, 0).text_frame.paragraphs[0].text
        if "Malang" in first_cell_text:
            set_paragraph_text(tbl.cell(0, 0).text_frame.paragraphs[0], f"Malang,  {TANGGAL_LAPORAN}")

    # ------------------------------------------------------------
    # SLIDE 29-31: LEMBAR BUKTI PERJALANAN RESES
    # Dokumen bukti perjalanan periode April-Mei tidak relevan untuk periode ini,
    # sehingga foto lama dihapus. Slide/placeholder tetap ada mengikuti template,
    # siap ditempeli dokumen bukti perjalanan yang sebenarnya oleh pengguna.
    # ------------------------------------------------------------
    for slide_idx in (29, 30):
        s = slides[slide_idx]
        for shape in list(s.shapes):
            if shape.shape_type == 13:
                shape._element.getparent().remove(shape._element)

    # ------------------------------------------------------------
    # SLIDE 32-51: LAMPIRAN FOTO/VISUAL PELAKSANAAN KEGIATAN RESES
    # ------------------------------------------------------------
    for i, keg in enumerate(KEGIATAN):
        s = slides[31 + i]
        content_pics = find_content_pictures(s)
        if len(content_pics) != len(keg["fotos"]):
            print(f'  [WARN] slide {32+i} (Kegiatan {keg["no"]}): '
                  f'{len(content_pics)} slot foto template vs {len(keg["fotos"])} foto data')
        for pic, foto_name in zip(content_pics, keg["fotos"]):
            replace_picture(pic, os.path.join(FOTO_DIR, foto_name))

    # ------------------------------------------------------------
    # SLIDE 52: PENUTUP BUKU - tidak diubah
    # ------------------------------------------------------------

    prs.save(OUTPUT_FILE)
    print(f"\n[OK] Buku Laporan Reses berhasil dibuat: {OUTPUT_FILE}")
    print(f"     Jumlah slide: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
